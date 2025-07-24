# parsers.py
import PyMuPDF  # PyMuPDF
import pptx
import docx
import pandas as pd
import io

def parse_file(file, chunk_size=500):
    file_name = file.name.lower()

    if file_name.endswith(".pdf"):
        text = parse_pdf(file)
    elif file_name.endswith(".pptx"):
        text = parse_pptx(file)
    elif file_name.endswith(".docx"):
        text = parse_docx(file)
    elif file_name.endswith(".csv"):
        text = parse_csv(file)
    elif file_name.endswith(".txt") or file_name.endswith(".md"):
        text = parse_txt(file)
    else:
        raise ValueError(f"Unsupported file format: {file_name}")
    
    return chunk_text(text, chunk_size)


def parse_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return text


def parse_pptx(file):
    prs = pptx.Presentation(io.BytesIO(file.read()))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def parse_docx(file):
    doc = docx.Document(io.BytesIO(file.read()))
    return "\n".join([para.text for para in doc.paragraphs])


def parse_csv(file):
    df = pd.read_csv(file)
    return df.to_string(index=False)


def parse_txt(file):
    return file.read().decode("utf-8")


def chunk_text(text, chunk_size):
    text = text.strip().replace("\n", " ").replace("\r", " ")
    chunks = []
    while len(text) > chunk_size:
        split_at = text.rfind(" ", 0, chunk_size)
        if split_at == -1:
            split_at = chunk_size
        chunks.append(text[:split_at].strip())
        text = text[split_at:].strip()
    if text:
        chunks.append(text)
    return chunks
